# SPDX-FileCopyrightText: Copyright (c) 2026 The Newton Developers
# SPDX-License-Identifier: Apache-2.0

"""
NVIDIA-style PowerPoint deck for Phase 3 benchmark — white background edition.

6 slides:
  1. Title
  2. The Task (shaft-Z plot + outcome boxes)
  3. Benchmark Design (identical vs must-differ)
  4. Video comparison (physical placeholder + VBD + MuJoCo)
  5. Key Results (contact force, contact area, frame-time)
  6. Roadmap (thermal + principled + surrogate multiphysics)

Usage:  python make_pptx.py
Output: slides/phase3_benchmark.pptx
"""

from __future__ import annotations

import csv
import pathlib

import cv2
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

HERE  = pathlib.Path(__file__).parent
BENCH = HERE / "bench_results"
VIDS  = HERE / "videos"
OUT   = HERE / "slides"

# ── Palette (light mode) ───────────────────────────────────────────────────────
_BG      = RGBColor(0xFF, 0xFF, 0xFF)
_BG2     = RGBColor(0xF5, 0xF5, 0xF5)
_BG3     = RGBColor(0xEB, 0xEB, 0xEB)
_GREEN   = RGBColor(0x76, 0xB9, 0x00)
_GREEN_T = RGBColor(0xE8, 0xF5, 0xCC)
_GREEN_D = RGBColor(0x45, 0x70, 0x00)
_BLACK   = RGBColor(0x12, 0x12, 0x12)
_DARK    = RGBColor(0x2D, 0x2D, 0x2D)
_MID     = RGBColor(0x55, 0x55, 0x55)
_LGRAY   = RGBColor(0x88, 0x88, 0x88)
_RULE    = RGBColor(0xD8, 0xD8, 0xD8)
_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
_BLUE    = RGBColor(0x15, 0x6E, 0xD6)
_BLUE_T  = RGBColor(0xE3, 0xF0, 0xFF)
_ORANGE  = RGBColor(0xC8, 0x48, 0x00)
_ORANGE_T= RGBColor(0xFD, 0xEE, 0xE3)
_PURPLE  = RGBColor(0x6A, 0x3D, 0x9A)
_PURPLE_T= RGBColor(0xF0, 0xE8, 0xFF)
_TEAL    = RGBColor(0x00, 0x7A, 0x87)
_TEAL_T  = RGBColor(0xD8, 0xF4, 0xF7)

# Slide: 13.33" × 7.5"
_SW = Inches(13.33)
_SH = Inches(7.50)
_VID_AR = 1944 / 1108   # actual video dimensions

# Contact stiffness — stored for force estimate plot
_KH_VBD = 5e6
_KH_MJ  = 2e9

# ── Plot colours ───────────────────────────────────────────────────────────────
_PLT_BG    = "#FFFFFF"
_PLT_AX    = "#FAFAFA"
_PLT_GRID  = "#E8E8E8"
_PLT_SPINE = "#CCCCCC"
_PLT_TICK  = "#555555"
_PLT_LABEL = "#444444"
_C_VBD     = "#1565C0"
_C_MJ      = "#C84800"

_PHASE_PLT = {
    "vertical_demo": "#F2F4FF", "approach":      "#EBF4FF",
    "sit_on_hub":    "#FFFBF0", "search_groove": "#FFF5EB",
    "engage":        "#F0FAEB", "push_in":       "#E8F7E8",
    "friction_demo": "#FFF1EB", "torsion_lock":  "#FFE8E8",
    "pull_out":      "#E8F9FC",
}
_PHASE_SHORT = {
    "vertical_demo": "vertical",  "approach": "approach",
    "sit_on_hub":    "sit on hub","search_groove": "search groove",
    "engage":        "engage",    "push_in": "push in",
    "friction_demo": "friction",  "torsion_lock": "torsion lock",
    "pull_out":      "pull out",
}


# ── Data ───────────────────────────────────────────────────────────────────────

def _load(solver: str) -> dict:
    rows = list(csv.DictReader((BENCH / f"bench_{solver}.csv").open()))
    g = lambda k, d=0.0: np.array([float(r.get(k, d)) for r in rows])
    return dict(
        solver=solver,
        t=g("t"), shaft_z=g("shaft_z_mm"), frame_ms=g("frame_ms"),
        phase=[r["phase"] for r in rows],
        contact_area=g("contact_area_mm2"),
        n_contacts=g("n_contacts"),
        max_depth=g("max_depth_mm"),
    )


def _smooth(a: np.ndarray, w: int = 9) -> np.ndarray:
    if len(a) < w:
        return a
    return np.convolve(a, np.ones(w) / w, "same")


# ── Plot helpers ───────────────────────────────────────────────────────────────

def _base_ax(ax, ylabel="", xlabel=""):
    ax.set_facecolor(_PLT_AX)
    ax.figure.patch.set_facecolor(_PLT_BG)
    for sp in ax.spines.values():
        sp.set_color(_PLT_SPINE)
    ax.tick_params(colors=_PLT_TICK, labelsize=12)
    ax.yaxis.label.set_color(_PLT_LABEL)
    ax.xaxis.label.set_color(_PLT_LABEL)
    if ylabel: ax.set_ylabel(ylabel, fontsize=14, labelpad=8)
    if xlabel: ax.set_xlabel(xlabel, fontsize=14, labelpad=8)
    ax.grid(axis="y", color=_PLT_GRID, lw=1.0)


def _bands(ax, t, phases, labels=True):
    prev = None; t0 = t[0]
    for i, ph in enumerate(phases):
        if ph != prev:
            if prev is not None:
                ax.axvspan(t0, t[i], color=_PHASE_PLT.get(prev, "#F8F8F8"),
                           alpha=1.0, lw=0, zorder=0)
                if labels:
                    mid = (t0 + t[i]) / 2
                    ax.text(mid, 1.015, _PHASE_SHORT.get(prev, prev),
                            ha="center", va="bottom", fontsize=8, color="#AAAAAA",
                            transform=ax.get_xaxis_transform(), zorder=5)
            prev = ph; t0 = t[i]
    if prev:
        ax.axvspan(t0, t[-1], color=_PHASE_PLT.get(prev, "#F8F8F8"),
                   alpha=1.0, lw=0, zorder=0)
        if labels:
            mid = (t0 + t[-1]) / 2
            ax.text(mid, 1.015, _PHASE_SHORT.get(prev, prev),
                    ha="center", va="bottom", fontsize=8, color="#AAAAAA",
                    transform=ax.get_xaxis_transform(), zorder=5)


def _save(fig, name: str) -> pathlib.Path:
    p = OUT / f"_p_{name}.png"
    fig.savefig(str(p), dpi=150, bbox_inches="tight",
                facecolor=fig.get_facecolor())
    plt.close(fig)
    return p


# ── Plots ─────────────────────────────────────────────────────────────────────

def plot_shaft_z(vbd: dict, mj: dict) -> pathlib.Path:
    fig, ax = plt.subplots(figsize=(10, 4.8))
    _base_ax(ax, ylabel="Shaft Z  (mm)", xlabel="Simulation time  (s)")
    _bands(ax, mj["t"], mj["phase"])

    ax.axhline( 22.5, color="#BBBBBB", lw=1.2, ls="--", alpha=0.9, zorder=3)
    ax.axhline(  0.0, color="#DDDDDD", lw=0.8, ls=":",  alpha=0.8, zorder=3)
    ax.axhline(-22.5, color="#BBBBBB", lw=1.2, ls="--", alpha=0.9, zorder=3)
    ax.fill_between([mj["t"][0], mj["t"][-1]], -22.5, 22.5,
                    color="#E8EEFF", alpha=0.5, zorder=1)
    ax.text(0.45,  22.5 + 2.5, "hub top",    color="#AAAAAA", fontsize=10)
    ax.text(0.45, -22.5 - 6.0, "hub bottom", color="#AAAAAA", fontsize=10)

    ax.plot(vbd["t"], vbd["shaft_z"], color=_C_VBD, lw=3.0,
            label="VBD (AVBD)  —  6.8 mm insertion", zorder=5)
    ax.plot(mj["t"],  mj["shaft_z"],  color=_C_MJ,  lw=3.0,
            label="MuJoCo  —  72.5 mm insertion",    zorder=5)

    vmin = int(np.argmin(vbd["shaft_z"]))
    ax.annotate("AVBD holds shaft\non tooth lands",
                xy=(vbd["t"][vmin], vbd["shaft_z"][vmin]),
                xytext=(vbd["t"][vmin] - 3.0, vbd["shaft_z"][vmin] - 20),
                color=_C_VBD, fontsize=11, ha="center",
                arrowprops=dict(arrowstyle="->", color=_C_VBD, lw=1.5), zorder=6)
    snap = int(np.argmin(np.diff(mj["shaft_z"])))
    ax.annotate(f"snap-to-groove\n@ {mj['t'][snap+1]:.1f} s",
                xy=(mj["t"][snap + 1], mj["shaft_z"][snap + 1]),
                xytext=(mj["t"][snap + 1] + 1.8, mj["shaft_z"][snap + 1] + 28),
                color=_C_MJ, fontsize=11, ha="center",
                arrowprops=dict(arrowstyle="->", color=_C_MJ, lw=1.5), zorder=6)

    ax.set_xlim(mj["t"][0], mj["t"][-1])
    ax.legend(fontsize=12, framealpha=0.95, facecolor="#FFFFFF",
              labelcolor="#222222", edgecolor="#CCCCCC", loc="lower right")
    fig.tight_layout(rect=[0, 0, 1, 0.92])
    return _save(fig, "shaft_z")


def plot_contact_force(vbd: dict, mj: dict) -> pathlib.Path:
    """Estimated contact force F = kh × max_depth × contact_area.

    Dual y-axes show the ~1000× scale difference between position-level
    (VBD AVBD, left axis) and force-level (MuJoCo, right axis) contacts.
    """
    f_vbd = _KH_VBD * (vbd["max_depth"] / 1000.0) * (vbd["contact_area"] / 1e6)
    f_mj  = _KH_MJ  * (mj["max_depth"]  / 1000.0) * (mj["contact_area"]  / 1e6)
    f_vbd_s = _smooth(f_vbd, 11)
    f_mj_s  = _smooth(f_mj,  11)

    fig, ax1 = plt.subplots(figsize=(10, 4.0))
    _base_ax(ax1, ylabel="", xlabel="Simulation time  (s)")
    _bands(ax1, mj["t"], mj["phase"])
    ax1.set_ylabel("VBD contact force  (N)", color=_C_VBD, fontsize=13, labelpad=8)
    ax1.tick_params(axis="y", colors=_C_VBD, labelsize=11)

    ax2 = ax1.twinx()
    ax2.set_facecolor("none")
    ax2.set_ylabel("MuJoCo contact force  (N)", color=_C_MJ, fontsize=13, labelpad=10)
    ax2.tick_params(axis="y", colors=_C_MJ, labelsize=11)
    for sp in ax2.spines.values(): sp.set_color(_PLT_SPINE)

    l1, = ax1.plot(vbd["t"], f_vbd_s, color=_C_VBD, lw=2.5,
                   label=f"VBD   peak {f_vbd.max():.0f} N", zorder=5)
    l2, = ax2.plot(mj["t"],  f_mj_s,  color=_C_MJ,  lw=2.5,
                   label=f"MuJoCo   peak {f_mj.max():.0f} N", zorder=5)

    ax1.set_ylim(bottom=0)
    ax2.set_ylim(bottom=0)
    ax1.set_xlim(mj["t"][0], mj["t"][-1])

    # Ratio annotation
    ratio = f_mj.max() / max(f_vbd.max(), 1)
    ax1.text(0.98, 0.96,
             f"Peak ratio  {ratio:.0f}×\n"
             f"kh differs  {int(_KH_MJ/_KH_VBD)}×",
             transform=ax1.transAxes, ha="right", va="top",
             fontsize=11, color="#555555",
             bbox=dict(boxstyle="round,pad=0.4", fc="#FFFFFF", ec="#CCCCCC"))

    lines = [l1, l2]
    labels = [l.get_label() for l in lines]
    ax1.legend(lines, labels, fontsize=11, framealpha=0.95, facecolor="#FFFFFF",
               labelcolor="#222222", edgecolor="#CCCCCC", loc="upper left")
    ax1.set_title("Estimated contact force  F = kh × δ × A", fontsize=13,
                  color="#555555", pad=10)
    fig.tight_layout(rect=[0, 0, 1, 0.92])
    return _save(fig, "contact_force")


def plot_contact_area(vbd: dict, mj: dict) -> pathlib.Path:
    fig, ax = plt.subplots(figsize=(6.5, 3.6))
    _base_ax(ax, ylabel="Contact area  (mm²)", xlabel="Simulation time  (s)")
    _bands(ax, mj["t"], mj["phase"], labels=False)
    ax.plot(vbd["t"], _smooth(vbd["contact_area"]), color=_C_VBD, lw=2.5,
            label=f"VBD   peak {vbd['contact_area'].max():.0f} mm²")
    ax.plot(mj["t"],  _smooth(mj["contact_area"]),  color=_C_MJ,  lw=2.5,
            label=f"MuJoCo   peak {mj['contact_area'].max():.0f} mm²")
    ax.set_xlim(mj["t"][0], mj["t"][-1])
    ax.set_ylim(bottom=0)
    ax.legend(fontsize=11, framealpha=0.95, facecolor="#FFFFFF",
              labelcolor="#222222", edgecolor="#CCCCCC")
    ax.set_title("Hydroelastic contact patch area", color="#555555", fontsize=12, pad=10)
    fig.tight_layout(rect=[0, 0, 1, 0.95])
    return _save(fig, "contact_area")


def plot_frame_time(vbd: dict, mj: dict) -> pathlib.Path:
    phases = list(dict.fromkeys(vbd["phase"] + mj["phase"]))
    vms = [np.mean([f for f, p in zip(vbd["frame_ms"], vbd["phase"]) if p == ph] or [0])
           for ph in phases]
    mms = [np.mean([f for f, p in zip(mj["frame_ms"],  mj["phase"])  if p == ph] or [0])
           for ph in phases]
    x, w = np.arange(len(phases)), 0.38
    labels = [_PHASE_SHORT.get(p, p).replace(" ", "\n") for p in phases]

    fig, ax = plt.subplots(figsize=(6.5, 3.6))
    ax.set_facecolor(_PLT_AX); fig.patch.set_facecolor(_PLT_BG)
    for sp in ax.spines.values(): sp.set_color(_PLT_SPINE)
    ax.tick_params(colors=_PLT_TICK, labelsize=9)
    ax.grid(axis="y", color=_PLT_GRID, lw=1.0)

    bv = ax.bar(x - w/2, vms, w, label="VBD",    color=_C_VBD, alpha=0.85)
    bm = ax.bar(x + w/2, mms, w, label="MuJoCo", color=_C_MJ,  alpha=0.85)
    ax.set_xticks(x); ax.set_xticklabels(labels, fontsize=8.5, color=_PLT_TICK)
    ax.set_ylabel("Mean frame time  (ms)", fontsize=12, color=_PLT_LABEL, labelpad=6)
    ax.set_ylim(bottom=0)
    ax.legend(fontsize=11, framealpha=0.95, facecolor="#FFFFFF",
              labelcolor="#222222", edgecolor="#CCCCCC")
    ax.set_title("Frame time per phase", color="#555555", fontsize=12, pad=10)
    for bar in list(bv) + list(bm):
        h = bar.get_height()
        if h > 2:
            c = _C_VBD if bar in bv else _C_MJ
            ax.text(bar.get_x() + bar.get_width() / 2, h + 1.5,
                    f"{h:.0f}", ha="center", fontsize=7.5, color=c)
    fig.tight_layout(rect=[0, 0, 1, 0.95])
    return _save(fig, "frame_time")


# ── Video thumbnail ────────────────────────────────────────────────────────────

def _thumb(video: pathlib.Path, t_sec: float = 9.0) -> pathlib.Path:
    out = OUT / f"_thumb_{video.stem}.png"
    cap = cv2.VideoCapture(str(video))
    cap.set(cv2.CAP_PROP_POS_MSEC, t_sec * 1000)
    ok, frame = cap.read(); cap.release()
    if ok:
        cv2.imwrite(str(out), frame)
    else:
        fig2, ax2 = plt.subplots(figsize=(4, 2.28))
        fig2.patch.set_facecolor("#E8E8E8"); ax2.set_facecolor("#E8E8E8"); ax2.axis("off")
        fig2.savefig(str(out), dpi=100, bbox_inches="tight",
                     facecolor=fig2.get_facecolor()); plt.close(fig2)
    return out


# ── PPTX helpers ───────────────────────────────────────────────────────────────

def _slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill; bg.solid(); bg.fore_color.rgb = _BG
    return s


def _box(s, l, t, w, h, fill, line=None, line_pt=0.75):
    sh = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(line_pt)
    return sh


def _tb(s, l, t, w, h):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    return tb, tf


def _run(tf, text, color=_DARK, size=14, bold=False, italic=False,
         align=PP_ALIGN.LEFT, space_before_pt=0.0, new_para=False,
         font="Calibri"):
    if new_para or tf.paragraphs[0].runs:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    p.alignment = align
    if space_before_pt > 0:
        p.space_before = Pt(space_before_pt)
    r = p.add_run()
    r.text = text; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = font


def _label(s, l, t, w, h, text, color=_DARK, size=14, bold=False,
           italic=False, align=PP_ALIGN.LEFT, font="Calibri"):
    _, tf = _tb(s, l, t, w, h)
    _run(tf, text, color=color, size=size, bold=bold, italic=italic,
         align=align, font=font)


def _img(s, path, l, t, w, h):
    s.shapes.add_picture(str(path), Inches(l), Inches(t), Inches(w), Inches(h))


def _movie(s, vid, poster, l, t, w, h):
    try:
        s.shapes.add_movie(str(vid), Inches(l), Inches(t), Inches(w), Inches(h),
                           str(poster), mime_type="video/mp4")
    except Exception:
        _img(s, poster, l, t, w, h)
        _label(s, l, t + h - 0.28, w, 0.25,
               f"[{vid.name}]", color=_LGRAY, size=9, align=PP_ALIGN.CENTER)


def _header(s, title, subtitle=""):
    _box(s, 0, 0, 13.33, 0.60, _GREEN)
    _label(s, 0.22, 0.10, 2.5, 0.40, "NVIDIA",
           color=_BG, size=22, bold=True, font="Calibri Light")
    _label(s, 8.80, 0.14, 4.35, 0.32, "Newton Physics Engine  ·  Phase 3",
           color=_BG, size=11, align=PP_ALIGN.RIGHT)
    if title:
        _label(s, 0.32, 0.68, 12.70, 0.60, title,
               color=_BLACK, size=28, bold=True, font="Calibri Light")
    if subtitle:
        _label(s, 0.32, 1.30, 12.70, 0.28, subtitle,
               color=_LGRAY, size=11.5, italic=True)
    _box(s, 0.32, 1.60, 12.70, 0.018, _RULE)


# ── Slide 1 – Title ────────────────────────────────────────────────────────────

def s1_title(prs):
    s = _slide(prs)
    _box(s, 0, 0, 0.18, 7.50, _GREEN)
    _label(s, 0.42, 0.28, 5.0, 0.60, "NVIDIA",
           color=_GREEN, size=32, bold=True, font="Calibri Light")
    _label(s, 0.42, 0.90, 8.0, 0.28,
           "Newton Physics Engine  ·  Phase 3  ·  Dynamics",
           color=_LGRAY, size=12)
    _label(s, 0.42, 1.38, 12.50, 1.0,
           "DIN 5480 Spline Insertion Dynamics",
           color=_BLACK, size=42, bold=True, font="Calibri Light")
    _label(s, 0.42, 2.40, 12.50, 0.55,
           "VBD  vs  MuJoCo  ·  Solver Benchmark",
           color=_GREEN, size=26, bold=False, font="Calibri Light")
    _box(s, 0.42, 3.10, 12.50, 0.025, _RULE)
    _label(s, 0.42, 3.26, 12.50, 0.28, "Core question",
           color=_LGRAY, size=12, italic=True)
    _label(s, 0.42, 3.58, 12.50, 0.70,
           "Which solver correctly captures tooth–groove engagement?\n"
           "What physical mechanisms explain the divergence?",
           color=_DARK, size=17, font="Calibri Light")
    chips = [("z = 6",     "DIN 5480 geometry"),
             ("23 s",      "scripted trajectory"),
             ("2 solvers", "VBD  /  MuJoCo"),
             ("Auto",      "headless benchmark")]
    chip_w = 2.88; x0 = 0.42
    for val, lbl in chips:
        _box(s, x0, 4.72, chip_w, 1.58, _BG2)
        _box(s, x0, 4.72, chip_w, 0.07, _GREEN)
        _label(s, x0 + 0.10, 5.00, chip_w - 0.20, 0.65, val,
               color=_GREEN_D, size=28, bold=True, align=PP_ALIGN.CENTER,
               font="Calibri Light")
        _label(s, x0 + 0.10, 5.68, chip_w - 0.20, 0.44, lbl,
               color=_MID, size=12, align=PP_ALIGN.CENTER)
        x0 += chip_w + 0.19


# ── Slide 2 – The Task ─────────────────────────────────────────────────────────

def s2_task(prs, sz_png):
    s = _slide(prs)
    _header(s, "The Task: Tooth–Groove Search and Engage",
            "DIN 5480  ·  z=6, module 4 mm, 30° pressure angle  ·  1 kg shaft  ·  40 mm engagement")
    _img(s, sz_png, 0.32, 1.70, 7.55, 3.62)
    rx = 8.10
    _label(s, rx, 1.72, 5.0, 0.34, "Why this is hard",
           color=_GREEN_D, size=15, bold=True)
    bullets = [
        "Teeth ride on land faces — near-zero insertion",
        "Groove alignment triggers a sudden force drop",
        "Solver must resolve that snap transition",
        "Contact geometry shifts every sub-step",
    ]
    y = 2.14
    for b in bullets:
        _, tf = _tb(s, rx + 0.08, y, 5.10, 0.40)
        _run(tf, f"▸  {b}", color=_MID, size=13); y += 0.44

    _box(s, rx, 3.82, 5.10, 1.62, _BLUE_T)
    _box(s, rx, 3.82, 0.08, 1.62, _BLUE)
    _label(s, rx + 0.20, 3.88, 4.80, 0.30, "VBD insertion depth",
           color=_BLUE, size=11.5, bold=True)
    _label(s, rx + 0.20, 4.18, 4.80, 0.70, "6.8 mm",
           color=_BLACK, size=40, bold=True, font="Calibri Light")
    _label(s, rx + 0.20, 4.90, 4.80, 0.28, "blocked by AVBD position correction",
           color=_MID, size=11.5)

    _box(s, rx, 5.56, 5.10, 1.62, _ORANGE_T)
    _box(s, rx, 5.56, 0.08, 1.62, _ORANGE)
    _label(s, rx + 0.20, 5.62, 4.80, 0.30, "MuJoCo insertion depth",
           color=_ORANGE, size=11.5, bold=True)
    _label(s, rx + 0.20, 5.92, 4.80, 0.70, "72.5 mm",
           color=_BLACK, size=40, bold=True, font="Calibri Light")
    _label(s, rx + 0.20, 6.64, 4.80, 0.28, "force-level contacts allow snap transition",
           color=_MID, size=11.5)
    _label(s, 0.32, 5.38, 7.55, 0.24,
           "Phase-coloured bands = scripted trajectory stages  ·  shaded region = hub bore",
           color=_LGRAY, size=9.5, italic=True)


# ── Slide 3 – Benchmark Design ─────────────────────────────────────────────────

def s3_design(prs):
    s = _slide(prs)
    _header(s, "Benchmark Design: What We Controlled",
            "Parameters held identical except where the contact model algorithm forces a difference")
    lx, rx = 0.32, 7.00
    col_w_l, col_w_r = 6.40, 6.01
    y0 = 1.78

    _box(s, lx, y0, col_w_l, 0.44, _GREEN_T)
    _box(s, lx, y0, 0.08, 0.44, _GREEN)
    _label(s, lx + 0.20, y0 + 0.06, col_w_l - 0.28, 0.34,
           "✓  Kept identical", color=_GREEN_D, size=14, bold=True)
    rows_id = [
        ("Geometry",         "DIN 5480 z=6, module 4 mm, 40 mm length"),
        ("Trajectory",       "Same 9-phase scripted insertion, 23 s"),
        ("Contact pipeline", "Newton hydroelastic SDF — both solvers"),
        ("Control gains",    "descent_k = 800 N·s/m  ·  μ = 0.30"),
        ("Warm-up",          "5 frames excluded from all timing"),
    ]
    y = y0 + 0.58
    for lbl, val in rows_id:
        _label(s, lx + 0.20, y, 2.60, 0.30, lbl + ":", color=_GREEN_D, size=13, bold=True)
        _label(s, lx + 2.84, y, 3.42, 0.30, val,        color=_MID,    size=13)
        _box(s, lx + 0.18, y + 0.34, col_w_l - 0.22, 0.012, _RULE)
        y += 0.94

    _box(s, rx, y0, col_w_r, 0.44, _ORANGE_T)
    _box(s, rx, y0, 0.08, 0.44, _ORANGE)
    _label(s, rx + 0.20, y0 + 0.06, col_w_r - 0.28, 0.34,
           "△  Must differ — algorithmic constraint", color=_ORANGE, size=14, bold=True)
    diffs = [
        ("Contact model",
         "VBD = position-level  (AVBD corrects penetration directly)",
         "MuJoCo = force-level  (contact force ∝ overlap depth)",
         "AVBD drives penetration to zero each sub-step; force model\n"
         "needs high stiffness to generate equivalent resistance"),
        ("Contact stiffness  kh",
         "VBD = 5 × 10⁶ N/m³",
         "MuJoCo = 2 × 10⁹ N/m³  (400× higher)",
         "Force-level equilibrium:  δ = F / (kh · A)\n"
         "400× stiffer → sub-mm penetration at operating loads"),
        ("Substeps + collision",
         "VBD = 8 substeps, collide once per frame",
         "MuJoCo = 16 substeps, collide every sub-step",
         "Tighter sub-dt + per-step geometry needed for\n"
         "force-level contacts on a rotating + descending shaft"),
    ]
    y = y0 + 0.58
    for hdr, v1, v2, reason in diffs:
        _label(s, rx + 0.20, y,        col_w_r - 0.28, 0.30, hdr,    color=_ORANGE, size=14, bold=True)
        _label(s, rx + 0.20, y + 0.34, col_w_r - 0.28, 0.26, v1,     color=_DARK,   size=12.5)
        _label(s, rx + 0.20, y + 0.62, col_w_r - 0.28, 0.26, v2,     color=_DARK,   size=12.5)
        _label(s, rx + 0.20, y + 0.92, col_w_r - 0.28, 0.44, reason, color=_LGRAY,  size=11, italic=True)
        _box(s, rx + 0.18, y + 1.40, col_w_r - 0.22, 0.012, _RULE)
        y += 1.58


# ── Slide 4 – Videos ───────────────────────────────────────────────────────────

def s4_videos(prs, vbd_vid, vbd_thumb, mj_vid, mj_thumb):
    s = _slide(prs)
    _header(s, "Side by Side: Physical Test vs Simulation",
            "Same insertion task — real hardware  |  VBD solver  |  MuJoCo solver")
    vid_w = 4.11; vid_h = vid_w / _VID_AR; y_vid = 1.72; gap = 0.22
    panels = [(_LGRAY, "Physical Test",     None,    None),
              (_BLUE,  "VBD Simulation",    vbd_vid, vbd_thumb),
              (_ORANGE,"MuJoCo Simulation", mj_vid,  mj_thumb)]
    for i, (color, lbl, vid, thumb) in enumerate(panels):
        xl = 0.28 + i * (vid_w + gap)
        _box(s, xl, y_vid, vid_w, 0.07, color)
        if vid is None:
            _box(s, xl, y_vid + 0.07, vid_w, vid_h, _BG2)
            _label(s, xl, y_vid + 0.07 + vid_h / 2 - 0.52, vid_w, 0.38,
                   "INSERT PHYSICAL\nTEST VIDEO HERE",
                   color=_LGRAY, size=14, bold=True, align=PP_ALIGN.CENTER,
                   font="Calibri Light")
            _label(s, xl, y_vid + 0.07 + vid_h / 2 + 0.14, vid_w, 0.26,
                   "(drag .mp4 here)",
                   color=_BG3, size=10, italic=True, align=PP_ALIGN.CENTER)
        else:
            _movie(s, vid, thumb, xl, y_vid + 0.07, vid_w, vid_h)
        _label(s, xl, y_vid + 0.07 + vid_h + 0.14, vid_w, 0.36,
               lbl, color=color, size=14, bold=True, align=PP_ALIGN.CENTER)
    y_div = y_vid + 0.07 + vid_h + 0.62
    _box(s, 0.28, y_div, 12.77, 0.018, _RULE)
    obs = [
        ("VBD:", _BLUE,
         "Shaft rotates in contact but never descends — "
         "AVBD position correction blocks every penetration attempt"),
        ("MuJoCo:", _ORANGE,
         "Shaft searches, then snaps into groove at 9.3 s — "
         "force-level contacts allow the engagement transition"),
    ]
    y = y_div + 0.22
    for prefix, color, text in obs:
        _, tf = _tb(s, 0.42, y, 12.50, 0.40)
        _run(tf, prefix + "  ", color=color, size=13, bold=True)
        p = tf.paragraphs[0]; r2 = p.add_run()
        r2.text = text; r2.font.size = Pt(13)
        r2.font.color.rgb = _MID; r2.font.name = "Calibri"
        y += 0.48


# ── Slide 5 – Key Results ──────────────────────────────────────────────────────

def s5_results(prs, cf_png, ca_png, ft_png):
    s = _slide(prs)
    _header(s, "Key Results: Contact Force and Solver Performance",
            "F = kh × δ × A  —  dual y-axes show the ~1000× contact force difference between position-level and force-level contacts")

    # Top: contact force — 9.5" wide, KPI boxes on right
    _img(s, cf_png, 0.32, 1.70, 9.72, 2.68)

    # KPI boxes — right column
    kpis = [
        (_BLUE,   _BLUE_T,   "VBD",    "6.8 mm",  "84 ms / frame",   "~13 N contact force"),
        (_ORANGE, _ORANGE_T, "MuJoCo", "72.5 mm", "114 ms / frame",  "~23 kN peak force"),
    ]
    y0 = 1.70
    for color, fill, solver, depth, speed, note in kpis:
        _box(s, 10.18, y0, 2.82, 1.32, fill)
        _box(s, 10.18, y0, 0.09, 1.32, color)
        _label(s, 10.34, y0 + 0.08, 2.55, 0.30,
               solver, color=color, size=13, bold=True)
        _label(s, 10.34, y0 + 0.38, 2.55, 0.54,
               depth, color=_BLACK, size=30, bold=True, font="Calibri Light")
        _label(s, 10.34, y0 + 0.94, 2.55, 0.22, speed, color=_MID,   size=11)
        _label(s, 10.34, y0 + 1.14, 2.55, 0.22, note,  color=_LGRAY, size=10.5, italic=True)
        y0 += 1.38

    # Bottom row
    _img(s, ca_png, 0.32, 4.48, 6.30, 2.88)
    _img(s, ft_png, 6.80, 4.48, 6.21, 2.88)


# ── Slide 6 – Roadmap ──────────────────────────────────────────────────────────

def s6_roadmap(prs):
    s = _slide(prs)
    _header(s, "Newton Multiphysics Roadmap",
            "From rigid-body dynamics to coupled thermal–mechanical simulation and learned surrogate models")

    # ── Phase timeline ──────────────────────────────────────────────────────
    phases_tl = [
        ("P1", "Geometry",   "DIN 5480\nparametric mesh", _GREEN,  True),
        ("P2", "Kinematic",  "SDF + hydro\ncontacts",     _GREEN,  True),
        ("P3", "Dynamics",   "VBD / MuJoCo\nbenchmark",   _GREEN,  True),
        ("P4", "Thermal",    "Friction heat\nWarp kernel", _ORANGE, False),
        ("P5", "Surrogate",  "PhysicsNeMo\ninference",     _PURPLE, False),
    ]
    box_w = 2.22; box_h = 1.08; gap = 0.22; y_tl = 1.78; x0 = 0.44

    for i, (tag, name, desc, color, done) in enumerate(phases_tl):
        xl = x0 + i * (box_w + gap)
        fill   = color if done else _BG2
        border = color
        txt_c  = _WHITE if done else color

        _box(s, xl, y_tl, box_w, box_h, fill, line=border, line_pt=1.5)
        _label(s, xl + 0.12, y_tl + 0.06, box_w - 0.24, 0.26,
               f"{tag}  {'✓' if done else '→'}",
               color=txt_c, size=12, bold=True)
        _label(s, xl + 0.12, y_tl + 0.32, box_w - 0.24, 0.24,
               name, color=txt_c, size=13, bold=True, font="Calibri Light")
        _label(s, xl + 0.12, y_tl + 0.58, box_w - 0.24, 0.42,
               desc, color=txt_c if done else _MID, size=10.5)

        # Arrow between boxes
        if i < len(phases_tl) - 1:
            ax_x = xl + box_w + gap * 0.12
            _label(s, ax_x, y_tl + box_h / 2 - 0.15, gap * 0.75, 0.30,
                   "→", color=_RULE, size=16, align=PP_ALIGN.CENTER)

    # ── Two detail columns ─────────────────────────────────────────────────
    y_det = y_tl + box_h + 0.28
    lx, rx = 0.44, 7.02
    col_w = 5.98; col_h = 4.30

    # Left — Phase 4: Thermal
    _box(s, lx, y_det, col_w, 0.44, _ORANGE_T)
    _box(s, lx, y_det, 0.09, 0.44, _ORANGE)
    _label(s, lx + 0.22, y_det + 0.07, col_w - 0.30, 0.30,
           "Phase 4  —  Thermal Simulation", color=_ORANGE, size=14, bold=True)
    _box(s, lx, y_det + 0.44, col_w, col_h - 0.44, _BG2)

    therm_rows = [
        ("Heat source",
         "Q = μ · N · v  from Phase 3 contact forces at each patch"),
        ("Warp kernel",
         "First-principles heat diffusion on shaft + hub surface mesh"),
        ("Thermal expansion",
         "δr = α · ΔT · r    (steel: α = 12 μm / m / °C)"),
        ("Effect on insertion",
         "50 °C rise → ~15 μm tighter clearance on z=6 shaft\n"
         "→ higher snap force, stiffer pull-out, shorter fatigue life"),
        ("Coupling direction",
         "Contact force  →  heat  →  expansion  →  clearance  →  force  (loop)"),
    ]
    y = y_det + 0.56
    for lbl, val in therm_rows:
        _label(s, lx + 0.22, y, 1.98, 0.26, lbl + ":", color=_ORANGE, size=11.5, bold=True)
        _label(s, lx + 2.22, y, col_w - 2.34, 0.40, val, color=_MID, size=11.5)
        _box(s, lx + 0.20, y + 0.48, col_w - 0.26, 0.010, _RULE)
        y += 0.70

    # Right — Phase 5: Surrogate
    _box(s, rx, y_det, col_w, 0.44, _PURPLE_T)
    _box(s, rx, y_det, 0.09, 0.44, _PURPLE)
    _label(s, rx + 0.22, y_det + 0.07, col_w - 0.30, 0.30,
           "Phase 5  —  Surrogate Multiphysics", color=_PURPLE, size=14, bold=True)
    _box(s, rx, y_det + 0.44, col_w, col_h - 0.44, _BG2)

    surr_items = [
        ("5a  Principled",
         _TEAL,
         "Coupled rigid-body + Warp heat diffusion kernel\n"
         "Full first-principles co-simulation at every frame"),
        ("5b  Surrogate",
         _PURPLE,
         "PhysicsNeMo model replaces contact + thermal kernel\n"
         "Trained on Phase 4 data; same trajectory driver interface"),
        ("Design goal",
         _DARK,
         "Swap physics backends without changing the simulation loop\n"
         "Plug-in point: predict(state) → force + temperature"),
        ("Expected gain",
         _GREEN_D,
         "5–10× wall-time reduction vs coupled co-simulation\n"
         "Enables real-time thermal-aware assembly planning"),
    ]
    y = y_det + 0.56
    for lbl, lbl_color, val in surr_items:
        _label(s, rx + 0.22, y, 2.00, 0.26, lbl + ":", color=lbl_color, size=11.5, bold=True)
        _label(s, rx + 2.24, y, col_w - 2.36, 0.44, val, color=_MID, size=11.5)
        _box(s, rx + 0.20, y + 0.54, col_w - 0.26, 0.010, _RULE)
        y += 0.76


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    OUT.mkdir(parents=True, exist_ok=True)

    print("Loading data ...")
    vbd = _load("vbd")
    mj  = _load("mujoco")

    print("Generating plots ...")
    sz_png = plot_shaft_z(vbd, mj)
    cf_png = plot_contact_force(vbd, mj)
    ca_png = plot_contact_area(vbd, mj)
    ft_png = plot_frame_time(vbd, mj)

    print("Extracting thumbnails ...")
    vbd_vid   = VIDS / "vbd_run.mp4"
    mj_vid    = VIDS / "mujoco_run.mp4"
    vbd_thumb = _thumb(vbd_vid,  t_sec=9.0)
    mj_thumb  = _thumb(mj_vid,   t_sec=9.0)

    print("Building PowerPoint ...")
    prs = Presentation()
    prs.slide_width  = _SW
    prs.slide_height = _SH

    s1_title(prs)
    s2_task(prs, sz_png)
    s3_design(prs)
    s4_videos(prs, vbd_vid, vbd_thumb, mj_vid, mj_thumb)
    s5_results(prs, cf_png, ca_png, ft_png)
    s6_roadmap(prs)

    out_path = OUT / "phase3_benchmark.pptx"
    prs.save(str(out_path))
    print(f"\nSaved -> {out_path}")


if __name__ == "__main__":
    main()
