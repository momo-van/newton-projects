# SPDX-FileCopyrightText: Copyright (c) 2026 The Newton Developers
# SPDX-License-Identifier: Apache-2.0
"""
Newton spline-insertion deck — v3 (sleek, large text, minimal boxes)

Usage:
    python slides/make_deck.py [--out slides/deck.pptx]
"""

from __future__ import annotations
import argparse, io, pathlib, tempfile
import cv2, numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt, Emu

# ── Paths ─────────────────────────────────────────────────────────────────────
HERE   = pathlib.Path(__file__).parent
ROOT   = HERE.parent
BENCH  = ROOT / "bench_results"
VIDS   = ROOT / "videos"
ASSETS = HERE.parent.parent   # spline_insertion/

# ── Newton palette ─────────────────────────────────────────────────────────────
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE   = RGBColor(0xF8, 0xF8, 0xF8)
CHARCOAL   = RGBColor(0x1A, 0x1A, 0x1A)
NAVY       = RGBColor(0x2B, 0x45, 0x90)
TEAL       = RGBColor(0x3A, 0x7A, 0x8C)
SAGE       = RGBColor(0x5E, 0xA8, 0x7A)
GREEN      = RGBColor(0x7D, 0xC8, 0x7A)
MINT       = RGBColor(0xA0, 0xD8, 0x7A)
LGREY      = RGBColor(0xF0, 0xF0, 0xF0)
MGREY      = RGBColor(0xCC, 0xCC, 0xCC)
DGREY      = RGBColor(0x88, 0x88, 0x88)
NAVY_DIM   = RGBColor(0x1C, 0x2E, 0x60)

# ── Canvas ────────────────────────────────────────────────────────────────────
W  = Inches(13.33)
H  = Inches(7.50)
ML = Inches(0.55)
MR = Inches(0.55)
MT = Inches(0.50)
CW = W - ML - MR   # ~12.23"


# ── Low-level primitives ──────────────────────────────────────────────────────

def rect(slide, l, t, w, h, fill=None, line=None, lw_pt=0.75):
    sh = slide.shapes.add_shape(1, int(l), int(t), int(w), int(h))
    sh.fill.solid() if fill else sh.fill.background()
    if fill: sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw_pt)
    return sh


def txt(slide, text, l, t, w, h,
        size=14, bold=False, italic=False,
        color=CHARCOAL, align=PP_ALIGN.LEFT,
        font="Calibri Light", wrap=True):
    tb = slide.shapes.add_textbox(int(l), int(t), int(w), int(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = font
    return tb


def hline(slide, l, t, w, color=MGREY, lw=0.75):
    ln = slide.shapes.add_connector(1, int(l), int(t), int(l+w), int(t))
    ln.line.color.rgb = color
    ln.line.width = Pt(lw)


def dots(slide, t, colors=None, r=Inches(0.14), gap=Inches(0.10), l=None):
    if colors is None:
        colors = [NAVY, TEAL, SAGE, GREEN, MINT]
    n = len(colors)
    total = n*(2*r) + (n-1)*gap
    x = (W - total) / 2 if l is None else l
    for c in colors:
        sh = slide.shapes.add_shape(9, int(x), int(t), int(2*r), int(2*r))
        sh.fill.solid(); sh.fill.fore_color.rgb = c
        sh.line.fill.background()
        x += 2*r + gap


def fit_pic(slide, path, l, t, w, h):
    im = Image.open(path)
    iw, ih = im.size
    s = min(w/iw, h/ih)
    sw, sh = int(iw*s), int(ih*s)
    ox = l + (w - Emu(sw))//2
    oy = t + (h - Emu(sh))//2
    slide.shapes.add_picture(str(path), int(ox), int(oy), Emu(sw), Emu(sh))


def thumb(vpath, t_sec=3.0):
    cap = cv2.VideoCapture(str(vpath))
    fps = cap.get(cv2.CAP_PROP_FPS) or 30
    cap.set(cv2.CAP_PROP_POS_FRAMES, int(t_sec*fps))
    ok, frm = cap.read(); cap.release()
    tmp = pathlib.Path(tempfile.mktemp(suffix=".png"))
    if ok: cv2.imwrite(str(tmp), frm)
    else:
        cv2.imwrite(str(tmp), np.full((360,640,3),40,dtype=np.uint8))
    return tmp


def video(slide, vpath, l, t, w, h, t_sec=3.0):
    th = thumb(vpath, t_sec)
    slide.shapes.add_movie(str(vpath), int(l), int(t), int(w), int(h),
                           poster_frame_image=str(th), mime_type="video/mp4")
    th.unlink(missing_ok=True)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── Compound components ───────────────────────────────────────────────────────

def slide_bg(slide, color=WHITE):
    rect(slide, 0, 0, W, H, fill=color)


def accent_bar(slide, color=NAVY, w=Inches(0.20)):
    """Thin vertical accent on left edge."""
    rect(slide, 0, 0, w, H, fill=color)


def slide_title_block(slide, title, subtitle=None,
                      title_size=34, sub_size=17,
                      title_color=NAVY, sub_color=TEAL,
                      l=None, t=None, w=None):
    l = l if l is not None else ML + Inches(0.25)
    t = t if t is not None else MT
    w = w if w is not None else CW - Inches(0.25)
    txt(slide, title, l, t, w, Inches(1.2),
        size=title_size, bold=True, color=title_color, font="Calibri Light")
    if subtitle:
        txt(slide, subtitle, l, t + Inches(0.90), w, Inches(0.55),
            size=sub_size, color=sub_color, font="Calibri Light")


def big_stat(slide, value, label, l, t, w=Inches(2.8),
             val_color=NAVY, val_size=38):
    txt(slide, value, l, t, w, Inches(0.90),
        size=val_size, bold=True, color=val_color,
        align=PP_ALIGN.CENTER, font="Calibri Light")
    txt(slide, label, l, t + Inches(0.82), w, Inches(0.35),
        size=11, color=DGREY, align=PP_ALIGN.CENTER, font="Calibri")


def bullet_list(slide, items, l, t, w, item_size=13.5, gap=Inches(0.42),
                dot_color=TEAL, text_color=CHARCOAL, bold_first=False):
    """items = list of str  or  (str, str) for (bold-head, body)."""
    y = t
    for item in items:
        # dot
        r = Inches(0.07)
        sh = slide.shapes.add_shape(9, int(l), int(y + Inches(0.08)),
                                    int(2*r), int(2*r))
        sh.fill.solid(); sh.fill.fore_color.rgb = dot_color
        sh.line.fill.background()
        if isinstance(item, tuple):
            head, body = item
            tb = slide.shapes.add_textbox(
                int(l + Inches(0.22)), int(y), int(w - Inches(0.22)), int(gap*1.6))
            tf = tb.text_frame; tf.word_wrap = True
            # head run
            p = tf.paragraphs[0]
            rh = p.add_run(); rh.text = head + "  "
            rh.font.bold = True; rh.font.size = Pt(item_size)
            rh.font.color.rgb = text_color; rh.font.name = "Calibri"
            rb = p.add_run(); rb.text = body
            rb.font.size = Pt(item_size - 0.5); rb.font.color.rgb = DGREY
            rb.font.name = "Calibri"
        else:
            txt(slide, item, l + Inches(0.22), y, w - Inches(0.22), gap,
                size=item_size, color=text_color, font="Calibri", wrap=True)
        y += gap
    return y


def section_label(slide, text, l, t, color=NAVY):
    txt(slide, text.upper(), l, t, CW, Inches(0.30),
        size=9, bold=True, color=color, font="Calibri",
        align=PP_ALIGN.LEFT)


def footer_dots(slide):
    dots(slide, H - Inches(0.32),
         colors=[NAVY, TEAL, SAGE, GREEN, MINT],
         r=Inches(0.09), gap=Inches(0.08), l=ML)


# ═══════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def s_title(prs):
    slide = blank(prs)
    slide_bg(slide)

    # Navy left bar
    rect(slide, 0, 0, Inches(0.30), H, fill=NAVY)

    # Decorative dot cluster (top right)
    dots(slide, Inches(0.35),
         colors=[MINT, GREEN, SAGE, TEAL, NAVY],
         r=Inches(0.19), gap=Inches(0.11),
         l=W - ML - 5*(Inches(0.38) + Inches(0.11)) + Inches(0.11))

    # Main title
    txt(slide, "DIN 5480 Spline Insertion", ML + Inches(0.40), Inches(1.60),
        Inches(9.0), Inches(1.30), size=46, bold=True,
        color=NAVY, font="Calibri Light")

    txt(slide, "VBD vs MuJoCo  ·  Solver Benchmark  ·  Newton GPU Physics",
        ML + Inches(0.40), Inches(2.85), Inches(10.0), Inches(0.55),
        size=20, color=TEAL, font="Calibri Light")

    hline(slide, ML + Inches(0.40), Inches(3.52), Inches(8.0), color=MGREY)

    # Stats row
    stats = [
        ("z = 6", "DIN 5480 spline geometry"),
        ("33 s",  "scripted trajectory"),
        ("2",     "solvers benchmarked"),
        ("RTX 5000", "Ada GPU  ·  CUDA"),
    ]
    sx = ML + Inches(0.40)
    sw = Inches(2.85)
    for v, l in stats:
        big_stat(slide, v, l, sx, Inches(3.65), w=sw, val_color=NAVY, val_size=30)
        sx += sw + Inches(0.10)

    footer_dots(slide)


def s_real_video(prs):
    slide = blank(prs)
    slide_bg(slide)
    accent_bar(slide, TEAL)

    section_label(slide, "The Real Process", ML + Inches(0.35), MT, TEAL)
    txt(slide, "Manual Spline Shaft Insertion",
        ML + Inches(0.35), MT + Inches(0.28), Inches(6.5), Inches(0.55),
        size=34, bold=True, color=NAVY, font="Calibri Light")
    txt(slide, "This is what Newton must simulate faithfully",
        ML + Inches(0.35), MT + Inches(0.85), Inches(7.0), Inches(0.35),
        size=15, color=TEAL, italic=True, font="Calibri Light")

    hline(slide, ML + Inches(0.35), Inches(1.55), CW - Inches(0.35), MGREY)

    # Portrait video on right (≈3" wide × 5.4" tall — 9:16 aspect)
    vid_w = Inches(3.10)
    vid_h = Inches(5.35)
    vid_l = W - MR - vid_w
    vid_t = Inches(1.70)
    video(slide, VIDS / "Assembly of spline shafts #shorts.mp4",
          vid_l, vid_t, vid_w, vid_h, t_sec=3.0)

    # Steps on left
    lw = vid_l - ML - Inches(0.55)
    steps = [
        (NAVY, "Approach",
         "Shaft lowered to hover just above the hub bore entrance."),
        (TEAL, "Angular search",
         "Operator applies light axial load and rotates slowly. "
         "Teeth ride on land faces — contact resistance is high."),
        (SAGE, "Snap-to-engage",
         "At tooth-groove alignment the axial resistance drops sharply. "
         "Shaft descends under gravity or applied push."),
        (GREEN, "Press-fit seat",
         "Operator pushes to full insertion depth. "
         "Groove walls resist torsion — the joint is locked."),
    ]
    y = Inches(1.72)
    for col, head, body in steps:
        sh = slide.shapes.add_shape(9, int(ML + Inches(0.35)),
                                    int(y + Inches(0.08)),
                                    int(Inches(0.18)), int(Inches(0.18)))
        sh.fill.solid(); sh.fill.fore_color.rgb = col; sh.line.fill.background()

        txt(slide, head, ML + Inches(0.65), y, lw - Inches(0.30), Inches(0.30),
            size=15, bold=True, color=col, font="Calibri")
        y += Inches(0.32)
        txt(slide, body, ML + Inches(0.65), y, lw - Inches(0.30), Inches(0.55),
            size=13, color=CHARCOAL, font="Calibri", wrap=True)
        y += Inches(0.72)

    txt(slide, "For z=6 there are 6 valid engagement angles (every 60°).\n"
               "SDF contact forces encode teeth-on-lands vs teeth-in-grooves\n"
               "— this is the core Newton demo moment.",
        ML + Inches(0.65), y + Inches(0.10), lw - Inches(0.30), Inches(0.80),
        size=12, color=DGREY, italic=True, font="Calibri", wrap=True)

    footer_dots(slide)


def s_geometry(prs):
    slide = blank(prs)
    slide_bg(slide)
    accent_bar(slide, SAGE)

    # Left: Rerun screenshot (large)
    img = ASSETS / "rerun_live.png"
    pic_w = Inches(6.8)
    fit_pic(slide, img, ML + Inches(0.25), MT, pic_w, H - MT - Inches(0.50))

    # Right panel
    rx = ML + Inches(0.25) + pic_w + Inches(0.45)
    rw = W - rx - MR

    section_label(slide, "Phase 1 — Geometry Creation", rx, MT, SAGE)
    slide_title_block(slide, "Parametric\nDIN 5480 Mesh",
                      subtitle="Involute spline  ·  z=6  ·  module=4 mm",
                      title_size=28, sub_size=14,
                      l=rx, t=MT + Inches(0.28), w=rw)

    hline(slide, rx, MT + Inches(1.78), rw, MGREY)

    params = [
        ("Teeth",           "z = 6  (DIN 5480)"),
        ("Module",          "m = 4 mm"),
        ("Pressure angle",  "α = 30°"),
        ("Spline length",   "L = 40 mm"),
        ("Hub outer dia.",  "OD = 64 mm"),
        ("Hub bore dia.",   "≈ 26 mm"),
    ]
    y = MT + Inches(1.92)
    for param, val in params:
        txt(slide, param, rx, y, rw * 0.50, Inches(0.30),
            size=12, color=DGREY, font="Calibri")
        txt(slide, val, rx + rw*0.50, y, rw*0.50, Inches(0.30),
            size=12, bold=True, color=CHARCOAL, font="Calibri")
        y += Inches(0.33)

    hline(slide, rx, y + Inches(0.05), rw, MGREY)
    y += Inches(0.22)

    txt(slide, "Generation pipeline", rx, y, rw, Inches(0.30),
        size=13, bold=True, color=TEAL, font="Calibri")
    y += Inches(0.32)

    pipeline = [
        "Shapely 2-D involute tooth profile",
        "Trimesh extrusion → watertight solid mesh",
        "USD export  (Z-up, metersPerUnit = 1)",
        "Newton SDF + hydroelastic contacts",
    ]
    bullet_list(slide, pipeline, rx, y, rw, item_size=12,
                gap=Inches(0.35), dot_color=SAGE)

    footer_dots(slide)


def s_kinematics(prs):
    slide = blank(prs)
    slide_bg(slide)
    accent_bar(slide, TEAL)

    section_label(slide, "Phase 2 — Kinematics", ML + Inches(0.35), MT, TEAL)
    slide_title_block(slide, "Search-and-Engage Insertion",
                      subtitle="Scripted trajectory  ·  8 phases  ·  ≤33 s total",
                      title_size=34, sub_size=16,
                      l=ML + Inches(0.35), t=MT + Inches(0.25))

    # Phase timeline bar
    phases = [
        ("vert.",      2.0,  NAVY),
        ("approach",   1.5,  TEAL),
        ("sit",        2.0,  SAGE),
        ("search",    20.0,  GREEN),
        ("hold",       0.5,  MINT),
        ("push",       2.5,  SAGE),
        ("retract",    1.0,  TEAL),
        ("withdraw",   3.5,  NAVY),
    ]
    total = sum(d for _,d,_ in phases)
    bar_l, bar_t = ML + Inches(0.35), Inches(2.05)
    bar_w, bar_h = CW - Inches(0.35), Inches(0.60)
    x = bar_l
    for ph, dur, col in phases:
        pw = int(bar_w * dur / total)
        rect(slide, x, bar_t, pw, bar_h, fill=col)
        if pw > Inches(0.55):
            txt(slide, ph, x + Inches(0.06), bar_t + Inches(0.08),
                pw - Inches(0.08), bar_h - Inches(0.12),
                size=9, bold=True, color=WHITE, font="Calibri")
        x += pw

    # Time ticks
    x = bar_l; cum = 0.0
    txt(slide, "0 s", x, bar_t + bar_h + Inches(0.04), Inches(0.4), Inches(0.25),
        size=8, color=DGREY, font="Calibri")
    for _, dur, _ in phases:
        cum += dur
        x += int(bar_w * dur / total)
        lbl = f"{cum:.0f} s" if cum < total else "≤33 s"
        txt(slide, lbl, x - Inches(0.20), bar_t + bar_h + Inches(0.04),
            Inches(0.45), Inches(0.25), size=8, color=DGREY, font="Calibri")

    # Two columns
    col_t = bar_t + bar_h + Inches(0.65)
    col_h = H - col_t - Inches(0.50)
    mid   = ML + Inches(0.35) + (CW - Inches(0.35)) / 2 + Inches(0.30)
    lw    = (CW - Inches(0.35)) / 2 - Inches(0.30)

    # Left: Real process steps
    txt(slide, "The Real Process", ML + Inches(0.35), col_t, lw, Inches(0.38),
        size=16, bold=True, color=NAVY, font="Calibri Light")

    process = [
        ("Approach", "Shaft lowered to just above hub bore entrance"),
        ("Angular search", "Light axial load + slow rotation; teeth ride on land faces"),
        ("Snap-to-engage", "Tooth-gap alignment → axial resistance drops, shaft descends"),
        ("Press-fit seat", "Pushed to full depth; groove walls resist torsion"),
    ]
    y = col_t + Inches(0.46)
    for i, (h, b) in enumerate(process):
        col = [NAVY, TEAL, SAGE, GREEN][i]
        sh = slide.shapes.add_shape(9, int(ML+Inches(0.35)), int(y+Inches(0.05)),
                                    int(Inches(0.16)), int(Inches(0.16)))
        sh.fill.solid(); sh.fill.fore_color.rgb = col; sh.line.fill.background()
        txt(slide, f"{h}  —  {b}",
            ML + Inches(0.62), y, lw - Inches(0.30), Inches(0.42),
            size=13, color=CHARCOAL, font="Calibri", wrap=True)
        y += Inches(0.44)

    # Right: Shaft control
    txt(slide, "Shaft Control", mid, col_t, lw, Inches(0.38),
        size=16, bold=True, color=NAVY, font="Calibri Light")

    ctrl = [
        ("Kinematic body", "Zero mass — position driven by Warp GPU kernel every substep"),
        ("Scripted trajectory", "Returns (ω₀, v_z, phase, done) from t and shaft_z_mm"),
        ("Event-driven done", "Groove catch fires _groove_caught_t; post-catch phases are relative, not wall-clock"),
        ("Smooth ramps", "Cubic ease-in/out on all transitions to avoid velocity spikes"),
    ]
    y = col_t + Inches(0.46)
    for h, b in ctrl:
        txt(slide, f"•  {h}  —  {b}", mid, y, lw, Inches(0.42),
            size=13, color=CHARCOAL, font="Calibri", wrap=True)
        y += Inches(0.44)

    # Vertical divider
    ln = slide.shapes.add_connector(1, int(mid - Inches(0.20)), int(col_t),
                                    int(mid - Inches(0.20)), int(H - Inches(0.45)))
    ln.line.color.rgb = MGREY; ln.line.width = Pt(0.75)

    footer_dots(slide)


def s_dynamics_design(prs):
    slide = blank(prs)
    slide_bg(slide)
    accent_bar(slide, NAVY)

    section_label(slide, "Phase 3 — Dynamics", ML + Inches(0.35), MT, NAVY)
    slide_title_block(slide, "Benchmark Design",
                      subtitle="Same scripted trajectory  ·  same geometry  ·  contact algorithm is the only variable",
                      title_size=34, sub_size=15,
                      l=ML + Inches(0.35), t=MT + Inches(0.25))

    hline(slide, ML + Inches(0.35), Inches(1.90), CW - Inches(0.35), MGREY)

    # Two columns
    mid  = W / 2 - Inches(0.10)
    lcol = ML + Inches(0.35)
    lw   = mid - lcol - Inches(0.25)
    rw   = W - mid - MR - Inches(0.10)

    col_t = Inches(2.05)

    # Column headers
    txt(slide, "✓ Kept Identical", lcol, col_t, lw, Inches(0.40),
        size=15, bold=True, color=SAGE, font="Calibri")
    txt(slide, "≠ Must Differ", mid, col_t, rw, Inches(0.40),
        size=15, bold=True, color=TEAL, font="Calibri")

    identical = [
        ("Geometry",         "DIN 5480  z=6, m=4 mm, 40 mm length"),
        ("Trajectory",       "Same scripted_trajectory.py  ·  8 phases  ·  ≤33 s"),
        ("Shaft mass",       "1.0 kg  ·  solid cylinder inertia"),
        ("Contact pipeline", "Newton hydroelastic SDF  (both solvers)"),
        ("Substeps / frame", "8  —  equal for fair wall-time comparison"),
        ("Descent force",    "15 N axial  ·  rotation 30 deg/s during search"),
        ("Warm-up",          "5 frames excluded from all timing stats"),
    ]
    different = [
        ("Contact model",
         "VBD = position-level (AVBD)\n"
         "MuJoCo = force-level (impulse)"),
        ("Stiffness kₕ",
         "VBD  5×10⁶  N/m  (displacement)\n"
         "MuJoCo  2×10⁹  N/m  (force)  → 400× higher"),
        ("Collision frequency",
         "VBD  once per frame\n"
         "MuJoCo  once per substep  (8× more often)"),
        ("Solver iterations",
         "VBD  32 iters  ·  work units 256\n"
         "MuJoCo  15 iters  ·  work units 120"),
    ]

    y = col_t + Inches(0.45)
    for param, val in identical:
        txt(slide, param, lcol, y, lw * 0.38, Inches(0.30),
            size=12, bold=True, color=CHARCOAL, font="Calibri")
        txt(slide, val, lcol + lw*0.38, y, lw*0.62, Inches(0.30),
            size=12, color=DGREY, font="Calibri")
        y += Inches(0.37)

    # Vertical divider
    ln = slide.shapes.add_connector(1, int(mid - Inches(0.12)), int(col_t),
                                    int(mid - Inches(0.12)), int(H - Inches(0.45)))
    ln.line.color.rgb = MGREY; ln.line.width = Pt(0.75)

    y = col_t + Inches(0.45)
    for param, val in different:
        txt(slide, param, mid, y, rw, Inches(0.28),
            size=12, bold=True, color=CHARCOAL, font="Calibri")
        y += Inches(0.28)
        txt(slide, val, mid + Inches(0.20), y, rw - Inches(0.20), Inches(0.55),
            size=11.5, color=DGREY, font="Calibri", wrap=True)
        y += Inches(0.60)

    footer_dots(slide)


def s_videos(prs):
    slide = blank(prs)
    slide_bg(slide)
    accent_bar(slide, NAVY)

    section_label(slide, "Phase 3 — Dynamics", ML + Inches(0.35), MT, NAVY)
    slide_title_block(slide, "VBD vs MuJoCo  —  Same Trajectory, Different Physics",
                      subtitle=None,
                      title_size=30, l=ML + Inches(0.35), t=MT + Inches(0.28))

    hline(slide, ML + Inches(0.35), Inches(1.65), CW - Inches(0.35), MGREY)

    vid_w = (CW - Inches(0.35) - Inches(0.35)) / 2
    vid_h = Inches(4.00)
    vid_t = Inches(1.78)

    # VBD video
    vl = ML + Inches(0.35)
    video(slide, VIDS/"vbd.mp4", vl, vid_t, vid_w, vid_h, t_sec=6.0)

    # MuJoCo video
    rl = vl + vid_w + Inches(0.35)
    video(slide, VIDS/"mujoco.mp4", rl, vid_t, vid_w, vid_h, t_sec=6.0)

    # Stats below each video
    stat_t = vid_t + vid_h + Inches(0.18)

    # VBD stats
    sx = vl
    sw = vid_w / 3
    for v, l in [("34.9 mm","insertion depth"),("71.6 ms","mean frame time"),("0.23×","real-time factor")]:
        big_stat(slide, v, l, sx, stat_t, w=sw, val_color=NAVY, val_size=22)
        sx += sw

    # MuJoCo stats
    sx = rl
    for v, l in [("2.5 mm","insertion depth"),("51.3 ms","mean frame time"),("0.32×","real-time factor")]:
        big_stat(slide, v, l, sx, stat_t, w=sw, val_color=TEAL, val_size=22)
        sx += sw

    # Solver labels above videos
    txt(slide, "VBD (AVBD)  —  position-level contacts",
        vl, vid_t - Inches(0.30), vid_w, Inches(0.28),
        size=13, bold=True, color=NAVY, font="Calibri")
    txt(slide, "MuJoCo  —  force-level contacts",
        rl, vid_t - Inches(0.30), vid_w, Inches(0.28),
        size=13, bold=True, color=TEAL, font="Calibri")

    footer_dots(slide)


def _full_plot_slide(prs, img_path, section, title, accent):
    slide = blank(prs)
    slide_bg(slide)
    accent_bar(slide, accent)

    section_label(slide, section, ML + Inches(0.35), MT, accent)
    txt(slide, title, ML + Inches(0.35), MT + Inches(0.26),
        CW - Inches(0.35), Inches(0.50),
        size=22, bold=True, color=NAVY, font="Calibri Light")

    hline(slide, ML + Inches(0.35), Inches(1.30), CW - Inches(0.35), MGREY)

    fit_pic(slide, img_path,
            ML + Inches(0.20), Inches(1.38),
            CW - Inches(0.20), H - Inches(1.38) - Inches(0.35))
    footer_dots(slide)


def s_traj_plot(prs):
    _full_plot_slide(prs,
        BENCH / "comparison_plot_light.png",
        "Benchmark Results",
        "Insertion Trajectory, Frame Timing & Angular Search",
        NAVY)


def s_force_plot(prs):
    _full_plot_slide(prs,
        BENCH / "comparison_forces_light.png",
        "Benchmark Results",
        "Force & Torque on Shaft Teeth  [d(mv)/dt decomposition]",
        TEAL)


def s_contact_plot(prs):
    _full_plot_slide(prs,
        BENCH / "comparison_contacts_light.png",
        "Benchmark Results",
        "Contact Dynamics & Convergence",
        SAGE)


def s_solver_tuning(prs):
    slide = blank(prs)
    slide_bg(slide)
    accent_bar(slide, GREEN)

    section_label(slide, "Solver Tuning", ML + Inches(0.35), MT, GREEN)
    slide_title_block(slide, "Rules of Thumb",
                      subtitle="Lessons from running VBD and MuJoCo on a contact-rich insertion task",
                      title_size=34, sub_size=15, sub_color=TEAL,
                      l=ML + Inches(0.35), t=MT + Inches(0.25))

    hline(slide, ML + Inches(0.35), Inches(1.90), CW - Inches(0.35), MGREY)

    # Two columns of insights
    col_t = Inches(2.05)
    col_h = H - col_t - Inches(0.50)
    mid   = W / 2 + Inches(0.10)
    lcol  = ML + Inches(0.35)
    lw    = mid - lcol - Inches(0.35)
    rw    = W - mid - MR

    left_rules = [
        ("Contact stiffness is not interchangeable",
         "VBD position-level: kₕ = 5×10⁶ N/m (displacement penalty). "
         "MuJoCo force-level: kₕ = 2×10⁹ N/m (impulse). "
         "The 400× difference is physically required, not a tuning choice."),
        ("njmax must cover 3× your contact budget",
         "Elliptic cone friction adds 3 constraint rows per contact pair. "
         "With 132 contacts, MuJoCo needs njmax ≥ 396 — we use 512. "
         "Undersizing causes silent nefc overflow and wrong dynamics."),
        ("Equal substeps, not equal iterations",
         "Fair comparison = same N_SUBSTEPS (8). Report work units "
         "(substeps × iters) alongside wall time. "
         "MuJoCo: 8×15 = 120 vs VBD: 8×32 = 256 — MuJoCo does less work per frame."),
        ("Always exclude warm-up frames",
         "Warp JIT compiles kernels on first use. The first 5 frames cost "
         "50–200 ms each. Exclude them; report only steady-state timing."),
    ]
    right_rules = [
        ("Search force vs stiffness balance",
         "At kₕ = 2×10⁹, a 15 N descent force equilibrates at ~2 mm bore penetration. "
         "Our groove detection requires 5 mm — so MuJoCo cannot engage. "
         "Either lower kₕ or increase search force to overcome this."),
        ("Collision frequency affects snap detection",
         "VBD collides once per frame; MuJoCo once per substep (8× more). "
         "But MuJoCo force-level contacts still cannot deliver the "
         "position-level snap that VBD achieves at groove alignment."),
        ("Torsion stiffness differs by design",
         "VBD position-level groove walls: finite stiffness → 20.95 °/s drift. "
         "MuJoCo force-level groove walls: effectively rigid → 0 °/s. "
         "Neither is wrong — they model different constraint philosophies."),
        ("RTF context: both solvers are batch-capable",
         "VBD 0.23×, MuJoCo 0.32× RTF on RTX 5000 Ada. "
         "Neither is real-time, but both are suitable for "
         "GPU-parallel batch rollouts in robot learning pipelines."),
    ]

    y = col_t
    for h, b in left_rules:
        txt(slide, h, lcol, y, lw, Inches(0.30),
            size=13, bold=True, color=CHARCOAL, font="Calibri")
        y += Inches(0.30)
        txt(slide, b, lcol + Inches(0.15), y, lw - Inches(0.15), Inches(0.58),
            size=11.5, color=DGREY, font="Calibri", wrap=True)
        y += Inches(0.66)

    ln = slide.shapes.add_connector(1, int(mid - Inches(0.15)), int(col_t),
                                    int(mid - Inches(0.15)), int(H - Inches(0.45)))
    ln.line.color.rgb = MGREY; ln.line.width = Pt(0.75)

    y = col_t
    for h, b in right_rules:
        txt(slide, h, mid, y, rw, Inches(0.30),
            size=13, bold=True, color=CHARCOAL, font="Calibri")
        y += Inches(0.30)
        txt(slide, b, mid + Inches(0.15), y, rw - Inches(0.15), Inches(0.58),
            size=11.5, color=DGREY, font="Calibri", wrap=True)
        y += Inches(0.66)

    footer_dots(slide)


def s_roadmap(prs):
    slide = blank(prs)
    slide_bg(slide)
    accent_bar(slide, SAGE)

    section_label(slide, "Multi-Physics Roadmap", ML + Inches(0.35), MT, SAGE)
    slide_title_block(slide, "Newton Spline Insertion  —  5-Phase Demo",
                      subtitle="From parametric geometry to PhysicsNeMo surrogate inference",
                      title_size=32, sub_size=15, sub_color=TEAL,
                      l=ML + Inches(0.35), t=MT + Inches(0.25))

    hline(slide, ML + Inches(0.35), Inches(1.90), CW - Inches(0.35), MGREY)

    phases = [
        ("1", "Geometry",   SAGE,  True,
         "Parametric DIN 5480 mesh",
         "Shapely 2-D → Trimesh extrude → USD\n"
         "z=6, module=4 mm, hydroelastic-ready"),
        ("2", "Kinematic",  TEAL,  True,
         "Rotation + translation drive",
         "SDF + hydroelastic contacts\n"
         "Interactive keyboard/slider control"),
        ("3", "Dynamics",   NAVY,  True,
         "VBD / MuJoCo solver benchmark",
         "Virtual force/torque manipulation\n"
         "Groove-snap detection, solver comparison"),
        ("4", "Thermal",    GREEN, False,
         "Friction heating on contact patches",
         "Custom Warp heat-diffusion kernel\n"
         "Q = μ·N·v from Phase 3 forces"),
        ("5", "Surrogate",  MINT,  False,
         "PhysicsNeMo contact-force model",
         "Trained on Phase 3 rollouts\n"
         "Per-phase surrogate inference plug-in"),
    ]

    # Horizontal flow diagram
    ph_t = Inches(2.10)
    ph_w = (CW - Inches(0.35)) / len(phases)

    for i, (num, name, col, done, short, detail) in enumerate(phases):
        cx = ML + Inches(0.35) + i * ph_w + ph_w / 2
        cy = ph_t + Inches(0.30)
        r  = Inches(0.36)

        # Circle
        sh = slide.shapes.add_shape(9, int(cx - r), int(cy - r),
                                    int(2*r), int(2*r))
        if done:
            sh.fill.solid(); sh.fill.fore_color.rgb = col
        else:
            sh.fill.background()
            sh.line.color.rgb = col; sh.line.width = Pt(2.5)
        if not done:
            sh.line.color.rgb = col

        # Phase number
        txt(slide, num, cx - r, cy - r - Inches(0.04), 2*r, 2*r,
            size=18, bold=True,
            color=WHITE if done else col,
            align=PP_ALIGN.CENTER, font="Calibri Light")

        # Connector arrow (except last)
        if i < len(phases) - 1:
            ax = cx + r + Inches(0.06)
            nx = cx + ph_w - r - Inches(0.06)
            ln = slide.shapes.add_connector(1, int(ax), int(cy), int(nx), int(cy))
            ln.line.color.rgb = MGREY; ln.line.width = Pt(1.5)

        # Phase name below circle
        txt(slide, name, cx - ph_w/2 + Inches(0.05), cy + r + Inches(0.10),
            ph_w - Inches(0.10), Inches(0.35),
            size=14, bold=True, color=col if done else DGREY,
            align=PP_ALIGN.CENTER, font="Calibri")

        # Short description
        txt(slide, short, cx - ph_w/2 + Inches(0.05), cy + r + Inches(0.45),
            ph_w - Inches(0.10), Inches(0.30),
            size=10.5, color=CHARCOAL,
            align=PP_ALIGN.CENTER, font="Calibri")

        # Done/Upcoming badge
        badge = "DONE" if done else "UPCOMING"
        bc    = SAGE if done else DGREY
        txt(slide, badge, cx - ph_w/2 + Inches(0.05), cy - r - Inches(0.35),
            ph_w - Inches(0.10), Inches(0.28),
            size=8, bold=True, color=bc,
            align=PP_ALIGN.CENTER, font="Calibri")

    # Detail cards for Phase 4 and 5
    card_t = ph_t + Inches(2.05)
    hline(slide, ML + Inches(0.35), card_t - Inches(0.10), CW - Inches(0.35), MGREY)

    for i, (num, name, col, done, short, detail) in enumerate(phases[3:], 3):
        cx   = ML + Inches(0.35) + i * ph_w
        card_w = ph_w * 2 - Inches(0.20) if i == 3 else ph_w * 2 - Inches(0.20)
        # Each phase gets half the remaining width
        cw = (CW - Inches(0.35)) / 2 - Inches(0.20)
        cl = ML + Inches(0.35) + (i - 3) * (cw + Inches(0.40))

        txt(slide, f"Phase {num}  —  {name}", cl, card_t, cw, Inches(0.32),
            size=14, bold=True, color=col, font="Calibri")
        txt(slide, detail, cl + Inches(0.10), card_t + Inches(0.35),
            cw - Inches(0.10), Inches(1.20),
            size=12.5, color=CHARCOAL, font="Calibri", wrap=True)

    footer_dots(slide)


def s_findings(prs):
    slide = blank(prs)
    slide_bg(slide)
    accent_bar(slide, NAVY)

    section_label(slide, "Summary", ML + Inches(0.35), MT, NAVY)
    slide_title_block(slide, "Key Findings",
                      subtitle="VBD position-level contacts succeed where MuJoCo force-level contacts cannot engage",
                      title_size=34, sub_size=15,
                      l=ML + Inches(0.35), t=MT + Inches(0.25))

    hline(slide, ML + Inches(0.35), Inches(1.88), CW - Inches(0.35), MGREY)

    rows = [
        ("Metric",            "VBD (AVBD)",                          "MuJoCo"),
        ("Insertion depth",   "34.9 mm  ✔  full engagement",    "2.5 mm  ✘  no groove snap"),
        ("Groove detection",  "YES  at t = 7.6 s",                  "NO  —  20 s timeout"),
        ("Mean frame time",   "71.6 ms / frame",                    "51.3 ms / frame  (1.40× faster)"),
        ("Real-time factor",  "0.23×",                          "0.32×"),
        ("Contact model",     "Position-level  kₕ = 5×10⁶",
                              "Force-level  kₕ = 2×10⁹"),
        ("Torsion drift",     "20.95 °/s  (pos-level stiffness limit)",
                              "0.00 °/s  (force-level → rigid)"),
        ("Physics insight",   "Hard pos-level contacts snap cleanly at alignment",
                              "15 N descent balanced by contact before bore entry"),
    ]

    tl = ML + Inches(0.35)
    tt = Inches(2.00)
    tw = CW - Inches(0.35)
    col_fs = [0.30, 0.35, 0.35]
    col_ws = [int(tw * f) for f in col_fs]
    rh = (H - tt - Inches(0.50)) / len(rows)

    for ri, row in enumerate(rows):
        x = tl
        for ci, cell in enumerate(row):
            is_hdr = ri == 0
            bg = NAVY if is_hdr else (LGREY if ri % 2 == 0 else WHITE)
            fc = WHITE if is_hdr else (CHARCOAL if ci == 0 else (NAVY if ci==1 else TEAL))
            rect(slide, x, int(tt + ri*rh), col_ws[ci], int(rh), fill=bg)
            txt(slide, cell,
                x + Inches(0.10), int(tt + ri*rh) + Inches(0.04),
                col_ws[ci] - Inches(0.14), int(rh) - Inches(0.06),
                size=12 if ri > 0 else 12, bold=(ci==0 or ri==0),
                color=fc, font="Calibri", wrap=True)
            x += col_ws[ci]

    footer_dots(slide)


# ── Build ─────────────────────────────────────────────────────────────────────

def build(out: pathlib.Path):
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    steps = [
        (s_title,          "Title"),
        (s_real_video,     "Real process video"),
        (s_geometry,       "Geometry"),
        (s_kinematics,     "Kinematics"),
        (s_dynamics_design,"Benchmark Design"),
        (s_videos,         "VBD vs MuJoCo videos"),
        (s_traj_plot,      "Trajectory benchmark"),
        (s_force_plot,     "Force & torque"),
        (s_contact_plot,   "Contact dynamics"),
        (s_solver_tuning,  "Solver tuning"),
        (s_roadmap,        "Roadmap"),
        (s_findings,       "Key findings"),
    ]
    for i, (fn, label) in enumerate(steps, 1):
        fn(prs)
        print(f"  {i:2d}/{len(steps)}  {label}")

    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    print(f"\nSaved -> {out}  ({out.stat().st_size//1024//1024} MB)")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--out", default=None)
    a  = ap.parse_args()
    out = pathlib.Path(a.out) if a.out else HERE / "spline_insertion_deck.pptx"
    build(out)
